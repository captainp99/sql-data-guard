"""Fill the A4I 4-quadrant template (sample.pptx) with content, preserving its design.

Clones the template's textboxes and replaces only the bullet text — the title styling,
colors, fonts, bullet glyphs, rounded cards and hero border all stay exactly as designed.

Usage: python fill_template.py   ->   writes one-pager.pptx
"""
from copy import deepcopy
from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATE = "template.pptx"
OUTPUT = "one-pager.pptx"

# ---- edit this block (matches the sample.html dummy data) ----
DECK_TITLE = "sql-data-guard"
DECK_SUBTITLE = "A4I Hackathon — 4-quadrant one-pager (sample)"

# keys MUST match the template's quadrant titles (paragraph 0 of each textbox)
QUADRANTS = {
    "Problem & opportunity": [
        "LLM-generated SQL can't run as prepared statements",
        "SQL injection & data leaks risk GDPR / CCPA fines",
        "DB permissions can't express row/column-level rules",
    ],
    "Hypothesis": [
        "Coding assistants can ship a safety layer faster",
        "Built with Claude Code + custom skills & MCP",
        "Engineering practices: tests-first, CI on every push",
        "Spec-driven via reusable SKILL.md contracts",
        "Tools: sqlglot, Flask, flasgger, python-pptx",
    ],
    "Use cases and results": [
        "Verifies & auto-rewrites unsafe queries pre-execution",
        "9 test suites: unit, joins, updates, REST, DuckDB",
        "4 integrations: library, REST, MCP wrapper, Dify",
        "Multi-dialect parsing (SQLite, Postgres, +)",
    ],
    "Value & Impact": [
        "Blocks data breaches before they reach the DB",
        "Drop-in for any LLM-to-SQL app — zero DB changes",
        "Reusable as pip package + Docker image",
        "Next: more dialects, policy templates, auth",
    ],
}
# --------------------------------------------------------------


def para_text(p):
    return "".join(r.text for r in p.runs)


def set_para_text(p_el, text):
    """Set a cloned <a:p>'s text into its first run; drop the other runs (keep formatting)."""
    runs = p_el.findall(qn("a:r"))
    first = runs[0]
    first.find(qn("a:t")).text = text
    for extra in runs[1:]:
        p_el.remove(extra)


def fill_textbox(tf, bullets):
    txBody = tf._txBody
    paras = tf.paragraphs
    title_p = paras[0]._p

    # find a non-empty bullet paragraph to use as the styling template
    tmpl = None
    for p in paras[1:]:
        if para_text(p).strip():
            tmpl = p._p
            break
    if tmpl is None:                      # no bullets in template box; nothing to clone
        return

    # remove every existing bullet paragraph (keep title + any empty spacer paragraphs)
    for p in paras[1:]:
        if para_text(p).strip():
            txBody.remove(p._p)

    # append new bullets cloned from the template bullet (preserves font/size/colour/glyph)
    for b in bullets:
        newp = deepcopy(tmpl)
        set_para_text(newp, b)
        txBody.append(newp)


def walk(shapes):
    for sh in shapes:
        if sh.shape_type == 6:            # group -> recurse
            yield from walk(sh.shapes)
        elif sh.has_text_frame:
            yield sh


def set_header(tf):
    runs = []
    for p in tf.paragraphs:
        runs.extend(p.runs)
    if not runs:
        return
    runs[0].text = DECK_TITLE             # first run = title line
    if len(runs) > 1:
        runs[-1].text = DECK_SUBTITLE     # last run = subtitle line (after the line break)


prs = Presentation(TEMPLATE)
slide = prs.slides[0]

for sh in walk(slide.shapes):
    head = para_text(sh.text_frame.paragraphs[0]).strip()
    if head in QUADRANTS:
        fill_textbox(sh.text_frame, QUADRANTS[head])
    elif head.startswith("A4I Hackathon") or sh.shape_id == 3:
        set_header(sh.text_frame)

prs.save(OUTPUT)
print(f"OK wrote {OUTPUT}")
